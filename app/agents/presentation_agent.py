from __future__ import annotations

from pathlib import Path
from zipfile import ZIP_DEFLATED, ZipFile
from xml.sax.saxutils import escape

from app.config.settings import PROJECT_DIR, get_settings
from app.models.schemas import WorkflowState
from app.tools.file_tools import ensure_dir, write_text


class PresentationAgent:
    name = "presentation_agent"

    def run(self, state: WorkflowState) -> WorkflowState:
        settings = get_settings()
        presentation_dir = ensure_dir(settings.output_dir / "presentations")
        md_path = presentation_dir / f"{state.run_id}_final_presentation.md"
        canonical_md = PROJECT_DIR / "presentation" / "final_presentation.md"
        content = self._slides(state)
        write_text(md_path, content)
        write_text(canonical_md, content)
        pptx_path = PROJECT_DIR / "presentation" / "final_presentation.pptx"
        if self._try_write_pptx(state, pptx_path):
            state.presentation_path = str(pptx_path)
        else:
            state.presentation_path = str(md_path)
        state.status = "completed"
        state.add_note(self.name, f"Generated presentation at {state.presentation_path}.")
        return state

    def _slide_content(self, state: WorkflowState) -> list[tuple[str, str]]:
        if not state.use_mock_llm:
            return self._research_slide_content(state)

        score = state.evaluation.total
        confidence = state.analysis.get("confidence_score", 0)
        return [
            ("Title", f"Agentic Research & Decision Intelligence Platform\nTopic: {state.query}"),
            ("Motivation", "Complex decisions require planning, evidence, critique, synthesis, and transparent uncertainty—especially in clinical-adjacent settings."),
            ("Problem Statement", "Transform a broad research topic into a professional evidence-grounded report and presentation with traceability and reproducibility."),
            ("Proposed System", "Ten specialized agents coordinate through typed WorkflowState, conditional graph edges, SQLite memory, and deterministic mock evidence."),
            ("Multi-Agent Architecture", "Streamlit UI and FastAPI backend feed a LangGraph-style workflow. Agents call search, citation, visualization, report, presentation, and memory tools."),
            ("LangGraph Workflow", "Planner → Research → Analyst → Critic → Fact-Check → Visualization → Report → Evaluator → Presentation → Memory."),
            ("Conditional Routing", "Critic loops back to Research when evidence is weak. Evaluator loops back to Report Writer when quality score Q is below threshold τ = 78."),
            ("Agent Roles", "Planner decomposes; Research retrieves; Analyst structures; Critic challenges; Fact-Checker grounds; Visualization explains; Writer synthesizes; Evaluator scores."),
            ("Implementation Stack", "Python, LangGraph-compatible orchestration, FastAPI, Streamlit, SQLite, Matplotlib, LaTeX, Pytest, Docker."),
            ("Mock Mode Design", "USE_MOCK_LLM=true runs without paid API keys. A local corpus of 21 verified references keeps outputs deterministic for grading."),
            ("Report Generation Pipeline", "Sources → claim checks → tables → figures → Markdown → LaTeX → PDF → presentation artifacts."),
            ("Mathematical Formulation", "Task decomposition T = {t₁,…,tₙ}; confidence C = Σwᵢsᵢ/Σwᵢ; quality Q = αF+βR+γS+δC_q; revision when Q < τ."),
            ("API Endpoints", "GET /health; POST /run; GET /runs/{run_id}; GET /runs/{run_id}/report; GET /runs/{run_id}/report.pdf; GET /runs/{run_id}/presentation; GET /runs/{run_id}/figures."),
            ("Evaluation Methodology", "Rubric: factuality, relevance, completeness, structure, citation quality, clarity, visual quality, reproducibility (artifact quality, not clinical validity)."),
            ("Results", f"Final score: {score}/100. Confidence: {confidence}/100. Sources: {len(state.sources)}. Figures: {len(state.figures)}. LaTeX PDF and PPTX generated."),
            (
                "User Interface",
                "Streamlit demo UI with hero header, sidebar controls, and tabs for Overview, Workflow, Agents, Results, Figures, and Downloads.\n\n"
                "![UI overview](../report/assets/ui/ui_home.png)",
            ),
            (
                "Agent Workflow Interface",
                "Workflow tab shows the ten-agent pipeline and conditional routing (Critic→Research, Evaluator→Report Writer).\n\n"
                "![Workflow tab](../report/assets/ui/ui_workflow.png)",
            ),
            (
                "Demo Results",
                f"Evaluation {score}/100; {len(state.figures)} figures; report and presentation generated.\n\n"
                f"![Evaluation results](../report/assets/results/demo_evaluation_score.png)\n\n"
                f"![Generated figures](../report/assets/results/demo_generated_figures.png)",
            ),
            ("Testing and Reproducibility", "pytest -v; python scripts/run_demo.py; python scripts/build_report_pdf.py; python scripts/capture_ui_screenshots.py; Docker Compose."),
            ("Limitations", "Static mock corpus; metadata-level fact checking; no clinical validation; human review required; not a medical device."),
            (
                "AI Usage and Development Transparency",
                "External AI tools (ChatGPT, Codex-style assistant, Cursor AI, NotebookLM) supported planning, implementation, documentation, and presentation drafting. "
                "The implemented system itself contains separate internal runtime agents. "
                "The student defined requirements, supervised development, reviewed outputs, tested the system, and prepared final submission. "
                "Mock mode and tests support reproducibility.",
            ),
            ("Future Work", "Live retrieval, vector stores, entailment-based fact checking, human-in-the-loop approval, security hardening, formal clinical evaluation."),
            ("Conclusion", "A reproducible multi-agent prototype that produces inspectable, evidence-grounded academic deliverables suitable for university evaluation."),
        ]

    def _research_slide_content(self, state: WorkflowState) -> list[tuple[str, str]]:
        score = state.evaluation.total
        confidence = state.analysis.get("confidence_score", 0)
        findings = state.analysis.get("key_findings", [])
        finding_text = "\n".join(f"- {item}" for item in findings[:5]) if isinstance(findings, list) and findings else "Key findings are available in the final report."
        supported = sum(1 for check in state.claim_checks if check.supported)
        return [
            ("Title", f"Dynamic Research Report\nTopic: {state.query}"),
            ("Research Question", f"The system analyzed this prompt as a topic-specific research task:\n{state.query}"),
            ("Workflow", "Planner → Research → Analyst → Critic → Fact-Checker → Visualization → Report Writer → Evaluator → Presentation → Memory."),
            ("Evidence Base", f"Retrieved sources: {len(state.sources)}. Confidence score: {confidence}/100."),
            ("Key Findings", finding_text),
            ("Claim Support", f"Supported claims: {supported}/{len(state.claim_checks)}. Unsupported claims are marked for caution in the PDF."),
            ("Evaluation", f"Final artifact score: {score}/100. The score reflects source coverage, claim support, tables, figures, and reproducibility."),
            (
                "Figures",
                "The report includes topic-specific evidence charts.\n\n![Evidence themes](../report/assets/evidence_by_theme.png)",
            ),
            ("Limitations", "This is a generated research draft. Human review is required before policy, financial, legal, health, or operational use."),
            ("Conclusion", "The final PDF is the primary output: a prompt-specific research paper grounded in retrieved evidence and claim checks."),
        ]

    def _slides(self, state: WorkflowState) -> str:
        blocks = [
            "---",
            "title: Agentic Research & Decision Intelligence Platform",
            "subtitle: Final Course Project",
            "---",
            "",
        ]
        for idx, (title, body) in enumerate(self._slide_content(state), start=1):
            blocks.extend([f"# {idx}. {title}", "", body, "", "Speaker notes: See final report for full technical detail.", "", "---", ""])
        return "\n".join(blocks).rstrip() + "\n"

    def _try_write_pptx(self, state: WorkflowState, path: Path) -> bool:
        try:
            from pptx import Presentation
        except Exception:
            self._write_minimal_pptx(state, path)
            return True
        from pptx.util import Inches

        image_map = {
            "User Interface": PROJECT_DIR / "report" / "assets" / "ui" / "ui_home.png",
            "Agent Workflow Interface": PROJECT_DIR / "report" / "assets" / "ui" / "ui_workflow.png",
            "Demo Results": PROJECT_DIR / "report" / "assets" / "results" / "demo_evaluation_score.png",
        }
        prs = Presentation()
        slides = self._slide_content(state)
        for title, body in slides:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = title
            clean_body = body.split("![")[0].strip()
            slide.placeholders[1].text = clean_body
            image_path = image_map.get(title)
            if image_path and image_path.exists():
                slide.shapes.add_picture(str(image_path), Inches(0.6), Inches(1.8), width=Inches(8.8))
        prs.save(path)
        return True

    def _write_minimal_pptx(self, state: WorkflowState, path: Path) -> None:
        slides = self._slide_content(state)
        path.parent.mkdir(parents=True, exist_ok=True)
        overrides = [
            '<Override PartName="/ppt/presentation.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml"/>',
            '<Override PartName="/docProps/core.xml" ContentType="application/vnd.openxmlformats-package.core-properties+xml"/>',
            '<Override PartName="/docProps/app.xml" ContentType="application/vnd.openxmlformats-officedocument.extended-properties+xml"/>',
        ]
        for idx in range(1, len(slides) + 1):
            overrides.append(f'<Override PartName="/ppt/slides/slide{idx}.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.slide+xml"/>')
        content_types = f"""<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">
  <Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>
  <Default Extension="xml" ContentType="application/xml"/>
  {''.join(overrides)}
</Types>"""
        root_rels = """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="ppt/presentation.xml"/>
  <Relationship Id="rId2" Type="http://schemas.openxmlformats.org/package/2006/relationships/metadata/core-properties" Target="docProps/core.xml"/>
  <Relationship Id="rId3" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/extended-properties" Target="docProps/app.xml"/>
</Relationships>"""
        slide_ids = "".join(f'<p:sldId id="{256 + idx}" r:id="rId{idx}"/>' for idx in range(1, len(slides) + 1))
        presentation_rels = "".join(
            f'<Relationship Id="rId{idx}" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide" Target="slides/slide{idx}.xml"/>'
            for idx in range(1, len(slides) + 1)
        )
        presentation = f"""<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:presentation xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
  <p:sldIdLst>{slide_ids}</p:sldIdLst>
  <p:sldSz cx="12192000" cy="6858000" type="wide"/>
  <p:notesSz cx="6858000" cy="9144000"/>
</p:presentation>"""
        app_xml = f"""<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Properties xmlns="http://schemas.openxmlformats.org/officeDocument/2006/extended-properties"><Application>Agentic Research Platform</Application><Slides>{len(slides)}</Slides></Properties>"""
        core_xml = """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<cp:coreProperties xmlns:cp="http://schemas.openxmlformats.org/package/2006/metadata/core-properties" xmlns:dc="http://purl.org/dc/elements/1.1/"><dc:title>Agentic Research Platform Presentation</dc:title></cp:coreProperties>"""
        with ZipFile(path, "w", compression=ZIP_DEFLATED) as pptx:
            pptx.writestr("[Content_Types].xml", content_types)
            pptx.writestr("_rels/.rels", root_rels)
            pptx.writestr("docProps/app.xml", app_xml)
            pptx.writestr("docProps/core.xml", core_xml)
            pptx.writestr("ppt/presentation.xml", presentation)
            pptx.writestr(
                "ppt/_rels/presentation.xml.rels",
                f'<?xml version="1.0" encoding="UTF-8" standalone="yes"?><Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">{presentation_rels}</Relationships>',
            )
            for idx, (title, body) in enumerate(slides, start=1):
                pptx.writestr(f"ppt/slides/slide{idx}.xml", self._slide_xml(title, body))

    def _slide_xml(self, title: str, body: str) -> str:
        title = escape(title)
        body = escape(body)
        return f"""<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:sld xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
  <p:cSld><p:spTree>
    <p:nvGrpSpPr><p:cNvPr id="1" name=""/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr><p:grpSpPr/>
    <p:sp><p:nvSpPr><p:cNvPr id="2" name="Title"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr><p:spPr><a:xfrm><a:off x="700000" y="650000"/><a:ext cx="10800000" cy="900000"/></a:xfrm></p:spPr><p:txBody><a:bodyPr/><a:lstStyle/><a:p><a:r><a:rPr lang="en-US" sz="3400" b="1"/><a:t>{title}</a:t></a:r></a:p></p:txBody></p:sp>
    <p:sp><p:nvSpPr><p:cNvPr id="3" name="Body"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr><p:spPr><a:xfrm><a:off x="900000" y="1900000"/><a:ext cx="10000000" cy="3000000"/></a:xfrm></p:spPr><p:txBody><a:bodyPr/><a:lstStyle/><a:p><a:r><a:rPr lang="en-US" sz="2200"/><a:t>{body}</a:t></a:r></a:p></p:txBody></p:sp>
  </p:spTree></p:cSld><p:clrMapOvr><a:masterClrMapping/></p:clrMapOvr>
</p:sld>"""
